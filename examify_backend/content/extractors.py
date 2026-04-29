import re

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

MAX_PDF_CHARS = 8000


def _is_page_number(line):
    if re.fullmatch(r"\d+", line):
        return True
    if re.fullmatch(r"page\s*\d+(\s*/\s*\d+)?", line, re.IGNORECASE):
        return True
    return False


def _clean_text(text, keep_newlines=False):
    if not text:
        return ""

    lines = [line.strip() for line in text.splitlines()]
    cleaned_lines = [line for line in lines if line and not _is_page_number(line)]

    if keep_newlines:
        normalized = [re.sub(r"\s+", " ", line).strip() for line in cleaned_lines]
        return "\n".join(normalized).strip()

    flattened = " ".join(cleaned_lines)
    return re.sub(r"\s+", " ", flattened).strip()


def _clean_multiline(text):
    if not text:
        return ""

    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    cleaned = []
    blank = False
    for line in lines:
        if not line:
            if not blank:
                cleaned.append("")
            blank = True
            continue
        blank = False
        if _is_page_number(line):
            continue
        cleaned.append(line)

    while cleaned and cleaned[0] == "":
        cleaned.pop(0)
    while cleaned and cleaned[-1] == "":
        cleaned.pop()

    return "\n".join(cleaned).strip()


def _truncate(text, limit):
    if not text or limit is None:
        return text
    return text[:limit]


def extract_from_pdf(file_path):
    reader = PdfReader(file_path)
    chunks = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text:
            chunks.append(page_text)

    cleaned = _clean_text("\n".join(chunks))
    return _truncate(cleaned, MAX_PDF_CHARS)


def extract_from_docx(file_path):
    doc = Document(file_path)
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
    joined = "\n".join(paragraphs)
    return _clean_text(joined, keep_newlines=True)


def extract_from_pptx(file_path):
    presentation = Presentation(file_path)
    slides_output = []

    for index, slide in enumerate(presentation.slides, start=1):
        title_text = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text.strip()

        content_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text.strip()
            if not shape_text:
                continue
            if title_text and shape_text == title_text:
                continue
            content_lines.append(shape_text)

        header = f"Slide {index}: {title_text}".strip()
        if content_lines:
            content = "\n".join(content_lines)
            slide_block = f"{header}\n{content}"
        else:
            slide_block = header

        slides_output.append(slide_block)

    return _clean_multiline("\n\n".join(slides_output))


def extract_from_video_url(url):
    return f"VIDEO_URL: {url} - Content extraction pending. Use topic mode instead."


def extract_from_topic(topic, subtopics):
    topic_text = (topic or "").strip()
    cleaned_subtopics = [str(item).strip() for item in (subtopics or []) if str(item).strip()]
    if not cleaned_subtopics:
        return f"Topic: {topic_text}\nSubtopics:\n"

    lines = "\n".join(f"- {item}" for item in cleaned_subtopics)
    return f"Topic: {topic_text}\nSubtopics:\n{lines}"
